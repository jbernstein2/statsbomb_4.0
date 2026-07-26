"""
pptx_generator.py

Turns the CSV files produced by stats_processor.py (a team-comparison
KPI table + per-player stat tables) into a formatted PowerPoint
(.pptx) report styled after a StatsBomb-style match recap deck.

Deck structure:
  1. Title slide
  2. For each game phase (Build-Up & Possession, Defending,
     Progression, Final Third / Red Zone, Crossing), in order:
       a. Phase divider/title slide
       b. Phase comparison slide - simple grouped bar chart of team
          totals for every KPI in that phase
       c. One slide per KPI in the phase - a team-totals bar chart
          plus best/worst performer callouts for each team
       d. Appendix - full player-by-player table for that phase's
          KPIs, grouped by team, with the highest value in each
          column highlighted green and the lowest highlighted red
       e. Any team visuals that belong to this phase (average
          position maps land in Build-Up & Possession, shot charts
          land in Final Third / Red Zone), grouped by team
  3. Supplementary KPI table (secondary stats not in the headline KPIs)
  4. Full-time summary slide
"""

import io
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

import stats_processor as sp


# ============================================================
# LAYOUT / DESIGN CONSTANTS
# ============================================================

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

BG_DARK = RGBColor(0x0E, 0x0E, 0x0E)
BG_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xED, 0xE6, 0xD6)
GREY_LABEL = RGBColor(0x9A, 0x9A, 0x9A)
MUTED_TEXT = RGBColor(0x6B, 0x6B, 0x6B)
DARK_TEXT = RGBColor(0x18, 0x18, 0x18)
ROW_ALT_BG = RGBColor(0xF2, 0xF0, 0xEB)

HIGHLIGHT_BEST = RGBColor(0xC8, 0xE8, 0xC6)   # light green
HIGHLIGHT_WORST = RGBColor(0xF4, 0xC7, 0xC3)  # light red

HEADER_FONT = "Cambria"
BODY_FONT = "Calibri"

# Automatic logo resolution search paths
DEFAULT_LOGO_PATHS = [
    Path("assets/Brooklyn_FC_logo.svg.webp"),
    Path("assets/brooklyn_fc_logo.png"),
    Path("assets/brooklyn_fc_logo.jpg"),
    Path("assets/logo.png"),
    Path("brooklyn_fc_logo.png"),
]


# ============================================================
# COLOR & FILE HELPERS
# ============================================================

def _hex_to_rgb(hex_color):
    h = hex_color.strip().lstrip("#")
    if len(h) != 6:
        raise ValueError(f"'{hex_color}' is not a valid hex color (expected 6 hex digits)")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _readable_text_color(rgb_color):
    luminance = 0.299 * rgb_color[0] + 0.587 * rgb_color[1] + 0.114 * rgb_color[2]
    return DARK_TEXT if luminance > 150 else RGBColor(0xFF, 0xFF, 0xFF)


def _find_default_logo() -> "Path | None":
    """Finds an existing logo file automatically from common paths."""
    for p in DEFAULT_LOGO_PATHS:
        if p.exists():
            return p
    return None


# ============================================================
# LOW-LEVEL SLIDE HELPERS
# ============================================================

def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_background(slide, prs, rgb_color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False, italic=False,
                  color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name=BODY_FONT,
                  anchor=None, letter_spaced=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    display_text = " ".join(text) if letter_spaced else text
    run.text = display_text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def _add_footer(slide, prs, left_text, right_text, on_dark=False):
    color = GREY_LABEL if on_dark else MUTED_TEXT
    _add_textbox(slide, Inches(0.5), Inches(SLIDE_HEIGHT_IN - 0.45), Inches(6.0), Inches(0.35),
                 left_text, size=10, color=color)
    _add_textbox(slide, Inches(SLIDE_WIDTH_IN - 6.5), Inches(SLIDE_HEIGHT_IN - 0.45), Inches(6.0), Inches(0.35),
                 right_text, size=10, color=color, align=PP_ALIGN.RIGHT)


def _add_monogram(slide, center_x, center_y, diameter, letter, fill_color, text_color):
    left = Emu(int(center_x - diameter / 2))
    top = Emu(int(center_y - diameter / 2))
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Emu(int(diameter)), Emu(int(diameter)))
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_color
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = letter
    run.font.size = Pt(int(diameter / Inches(1) * 30))
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = HEADER_FONT
    return circle


def _add_logo(slide, center_x, center_y, max_dim, logo_path):
    """
    Adds a logo image centered at (center_x, center_y), aspect-scaled within max_dim.
    Converts image formats (like WebP) into an in-memory PNG stream for compatibility.
    """
    with Image.open(logo_path) as im:
        img_w, img_h = im.size

        # Convert image to RGBA PNG in memory to ensure full PPTX engine compatibility
        img_bytes = io.BytesIO()
        im.convert("RGBA").save(img_bytes, format="PNG")
        img_bytes.seek(0)

    aspect = img_w / img_h
    if aspect >= 1:
        w = max_dim
        h = Emu(int(w / aspect))
    else:
        h = max_dim
        w = Emu(int(h * aspect))

    left = Emu(int(center_x - w / 2))
    top = Emu(int(center_y - h / 2))

    return slide.shapes.add_picture(img_bytes, left, top, width=w, height=h)


def _add_decorative_arc(slide, prs):
    diameter = Inches(3.2)
    left = Inches(-1.6)
    top = Inches(SLIDE_HEIGHT_IN - 1.6)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    circle.fill.background()
    circle.line.color.rgb = RGBColor(0x3A, 0x3A, 0x3A)
    circle.line.width = Pt(1)
    circle.shadow.inherit = False
    return circle


def _style_table(table, header_color, header_text_color):
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = header_text_color
                r.font.size = Pt(14)
                r.font.name = BODY_FONT

    for row_idx in range(1, len(table.rows)):
        bg = ROW_ALT_BG if row_idx % 2 == 0 else BG_LIGHT
        for cell in table.rows[row_idx].cells:
            if cell.fill.type is None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13)
                    r.font.color.rgb = DARK_TEXT
                    r.font.name = BODY_FONT


def _add_grouped_bar_chart(slide, x, y, cx, cy, categories, team_name, team_color,
                            team_values, opponent_name, opponent_color, opponent_values):
    """Simple side-by-side team-vs-team bar chart for a list of KPIs."""
    chart_data = CategoryChartData()
    chart_data.categories = list(reversed(categories))
    chart_data.add_series(team_name, list(reversed(team_values)))
    chart_data.add_series(opponent_name, list(reversed(opponent_values)))

    gframe = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)
    chart = gframe.chart
    chart.has_legend = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(26)
    plot.data_labels.number_format = "0.##"
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.gap_width = 60
    plot.overlap = -10

    series = plot.series
    series[0].format.fill.solid()
    series[0].format.fill.fore_color.rgb = team_color
    series[0].format.line.fill.background()
    series[1].format.fill.solid()
    series[1].format.fill.fore_color.rgb = opponent_color
    series[1].format.line.fill.background()

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(11)
    category_axis.format.line.color.rgb = RGBColor(0xD5, 0xD5, 0xD5)

    value_axis = chart.value_axis
    value_axis.visible = False
    value_axis.has_major_gridlines = False
    value_axis.minimum_scale = 0

    return chart


def _add_team_legend(slide, team_name, team_color, opponent_name, opponent_color,
                      legend_x=Inches(7.6), legend_y=Inches(0.5)):
    chip1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_x, legend_y, Inches(0.28), Inches(0.28))
    chip1.fill.solid()
    chip1.fill.fore_color.rgb = team_color
    chip1.line.fill.background()
    chip1.shadow.inherit = False
    _add_textbox(slide, legend_x + Inches(0.38), legend_y - Inches(0.06), Inches(2.3), Inches(0.4),
                 team_name, size=13, bold=True, color=DARK_TEXT)

    chip2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_x + Inches(2.6), legend_y, Inches(0.28), Inches(0.28))
    chip2.fill.solid()
    chip2.fill.fore_color.rgb = opponent_color
    chip2.line.fill.background()
    chip2.shadow.inherit = False
    _add_textbox(slide, legend_x + Inches(2.98), legend_y - Inches(0.06), Inches(2.5), Inches(0.4),
                 opponent_name, size=13, bold=True, color=DARK_TEXT)


# ============================================================
# 1. TITLE SLIDE
# ============================================================

def add_title_slide(prs, team_name, team_color, opponent_name, opponent_color,
                     kicker="STATSBOMB MATCH RECAP",
                     subtitle="STATISTICAL MATCH ANALYSIS  \u2022  PLAYER & TEAM PERFORMANCE BREAKDOWN",
                     date_label=None, logo_image=None):
    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_DARK)
    _add_decorative_arc(slide, prs)

    _add_textbox(slide, Inches(0.8), Inches(1.55), Inches(9.0), Inches(0.4),
                 kicker, size=13, color=GREY_LABEL, letter_spaced=True)

    _add_textbox(slide, Inches(0.75), Inches(2.0), Inches(9.0), Inches(1.0),
                 team_name.upper(), size=44, bold=True, color=team_color, font_name=HEADER_FONT)

    _add_textbox(slide, Inches(0.8), Inches(2.95), Inches(9.0), Inches(0.7),
                 opponent_name.upper(), size=28, italic=True, color=CREAM, font_name=HEADER_FONT)

    _add_textbox(slide, Inches(0.8), Inches(3.85), Inches(9.5), Inches(0.4),
                 subtitle, size=12, color=GREY_LABEL, letter_spaced=False)

    footer_right = f"{team_name.upper()}   vs   {opponent_name.upper()}"
    _add_footer(slide, prs, date_label or "", footer_right, on_dark=True)

    return slide


# ============================================================
# 2. PHASE DIVIDER / TITLE SLIDE
# ============================================================

def add_phase_title_slide(prs, phase_index, phase_count, phase_name, metrics,
                           team_name, opponent_name, team_color, opponent_color,
                           logo_image=None):
    """Dark section-break slide introducing the upcoming game phase."""
    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_DARK)
    _add_decorative_arc(slide, prs)

    _add_textbox(slide, Inches(0.9), Inches(1.7), Inches(8.0), Inches(0.4),
                 f"PHASE {phase_index} OF {phase_count}", size=13, color=GREY_LABEL, letter_spaced=True)

    _add_textbox(slide, Inches(0.85), Inches(2.2), Inches(9.0), Inches(1.3),
                 phase_name.upper(), size=42, bold=True, color=CREAM, font_name=HEADER_FONT)

    kpi_list = " \u2022 ".join(label for label, _, _ in metrics)
    _add_textbox(slide, Inches(0.9), Inches(3.55), Inches(9.5), Inches(1.2),
                 kpi_list, size=13, color=GREY_LABEL)

    # Logo, top-right corner (falls back to a monogram if no logo is found)
    logo_to_use = logo_image if (logo_image and Path(logo_image).exists()) else _find_default_logo()

    if logo_to_use is not None and Path(logo_to_use).exists():
        _add_logo(slide, Inches(12.2), Inches(1.05), Inches(1.3), logo_to_use)
    else:
        monogram_letter = team_name.strip()[0].upper() if team_name.strip() else "?"
        _add_monogram(slide, Inches(12.2), Inches(1.05), Inches(1.3), monogram_letter, CREAM, BG_DARK)

    footer_right = f"{team_name.upper()}   vs   {opponent_name.upper()}"
    _add_footer(slide, prs, "Game Phase", footer_right, on_dark=True)

    return slide


# ============================================================
# 3. PHASE COMPARISON SLIDE (team totals bar chart, scoped to a phase)
# ============================================================

def add_phase_comparison_slide(prs, phase_name, comparison_df, metrics,
                                team_name, opponent_name, team_color, opponent_color):
    """Grouped bar chart of team totals for every KPI in a single game phase."""
    labels = [label for label, _, _ in metrics]

    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_LIGHT)

    _add_textbox(slide, Inches(0.5), Inches(0.28), Inches(6.5), Inches(0.35),
                 phase_name.upper(), size=13, bold=True, color=MUTED_TEXT, letter_spaced=True)
    _add_textbox(slide, Inches(0.5), Inches(0.6), Inches(6.0), Inches(0.55),
                 "PHASE TOTALS", size=24, bold=True, color=DARK_TEXT, font_name=HEADER_FONT)

    _add_team_legend(slide, team_name, team_color, opponent_name, opponent_color)

    rows = comparison_df.set_index("Metric").reindex(labels).reset_index()

    _add_grouped_bar_chart(
        slide, Inches(0.6), Inches(1.35), Inches(SLIDE_WIDTH_IN - 1.2), Inches(5.6),
        rows["Metric"].tolist(), team_name, team_color, rows[team_name].fillna(0).tolist(),
        opponent_name, opponent_color, rows[opponent_name].fillna(0).tolist(),
    )

    _add_footer(slide, prs, f"{phase_name} \u2014 Team Totals", f"{team_name.upper()} vs {opponent_name.upper()}")
    return slide


# ============================================================
# 4. PER-METRIC SLIDE (team totals bar chart + best/worst performer callouts)
# ============================================================

def add_metric_slide(prs, phase_name, label, player_col, team_value, opponent_value,
                      team_players_df, opponent_players_df,
                      team_name, opponent_name, team_color, opponent_color,
                      direction=None, include_goalkeepers=True,
                      stat_index=None, stat_total=None):
    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_LIGHT)

    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(6.5), Inches(0.3),
                 phase_name.upper(), size=11, bold=True, color=MUTED_TEXT, letter_spaced=True)
    _add_textbox(slide, Inches(0.5), Inches(0.48), Inches(8.3), Inches(0.55),
                 label.upper(), size=24, bold=True, color=DARK_TEXT, font_name=HEADER_FONT)

    if stat_index is not None and stat_total is not None:
        _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(4.0), Inches(0.28),
                     f"Statistic {stat_index} of {stat_total}", size=10, color=MUTED_TEXT)

    _add_textbox(slide, Inches(9.2), Inches(0.32), Inches(3.6), Inches(0.3),
                 "TEAM TOTAL", size=11, color=MUTED_TEXT, align=PP_ALIGN.RIGHT)

    def _fmt(v):
        if v is None or pd.isna(v):
            return "-"
        return str(int(v)) if float(v).is_integer() else f"{v:.2f}"

    _add_textbox(slide, Inches(9.0), Inches(0.55), Inches(1.8), Inches(0.75),
                 _fmt(team_value), size=32, bold=True, color=team_color, align=PP_ALIGN.CENTER,
                 font_name=HEADER_FONT)
    _add_textbox(slide, Inches(10.8), Inches(0.55), Inches(1.9), Inches(0.75),
                 _fmt(opponent_value), size=32, bold=True, color=opponent_color, align=PP_ALIGN.CENTER,
                 font_name=HEADER_FONT)
    _add_textbox(slide, Inches(9.0), Inches(1.25), Inches(1.8), Inches(0.3),
                 team_name, size=10, color=MUTED_TEXT, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(10.8), Inches(1.25), Inches(1.9), Inches(0.3),
                 opponent_name, size=10, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

    # Team totals bar chart for this single stat
    chart_top = Inches(1.9)
    chart_h = Inches(2.1)
    chart_left = Inches(2.2)
    chart_width = Inches(SLIDE_WIDTH_IN - 4.4)

    _add_grouped_bar_chart(
        slide, chart_left, chart_top, chart_width, chart_h,
        [label], team_name, team_color, [team_value if team_value is not None else 0],
        opponent_name, opponent_color, [opponent_value if opponent_value is not None else 0],
    )

    # Best / worst performer callouts, grouped by team
    callout_top = Inches(4.2)
    half_w = Inches((SLIDE_WIDTH_IN - 1.4) / 2)
    left_left = Inches(0.6)
    right_left = Inches(0.6) + half_w + Inches(0.2)

    effective_direction = direction or "higher"

    for side_left, df, name, color in (
        (left_left, team_players_df, team_name, team_color),
        (right_left, opponent_players_df, opponent_name, opponent_color),
    ):
        _add_textbox(slide, side_left, callout_top, half_w, Inches(0.35),
                     f"{name.upper()} \u2014 TOP PERFORMERS", size=14, bold=True, color=DARK_TEXT)

        (best_player, best_val), (worst_player, worst_val) = sp.best_worst_player(
            df, player_col, direction=effective_direction, include_goalkeepers=include_goalkeepers,
        )

        def _fmt_val(v):
            if v is None or pd.isna(v):
                return ""
            return str(int(v)) if float(v).is_integer() else f"{v:.2f}"

        if best_player is None:
            _add_textbox(slide, side_left, callout_top + Inches(0.45), half_w, Inches(0.4),
                         "No player data available", size=13, color=MUTED_TEXT)
            continue

        card_w = Inches((half_w - Inches(0.3)) / 2)

        _add_textbox(slide, side_left, callout_top + Inches(0.5), card_w, Inches(0.32),
                     "BEST", size=16, bold=True, color=color)
        _add_textbox(slide, side_left, callout_top + Inches(0.85), card_w, Inches(0.7),
                     str(best_player), size=24, bold=True, color=DARK_TEXT)
        _add_textbox(slide, side_left, callout_top + Inches(1.58), card_w, Inches(0.95),
                     _fmt_val(best_val), size=30, bold=True, color=color, font_name=HEADER_FONT)

        worst_left = side_left + card_w + Inches(0.3)
        _add_textbox(slide, worst_left, callout_top + Inches(0.5), card_w, Inches(0.32),
                     "WORST", size=16, bold=True, color=MUTED_TEXT)
        _add_textbox(slide, worst_left, callout_top + Inches(0.85), card_w, Inches(0.7),
                     str(worst_player), size=24, bold=True, color=DARK_TEXT)
        _add_textbox(slide, worst_left, callout_top + Inches(1.58), card_w, Inches(0.95),
                     _fmt_val(worst_val), size=30, bold=True, color=MUTED_TEXT, font_name=HEADER_FONT)

    footer_left = f"{phase_name} \u2014 {label}"
    footer_right = f"{team_name.upper()} vs {opponent_name.upper()}"
    _add_footer(slide, prs, footer_left, footer_right)
    return slide


# ============================================================
# 5. APPENDIX - FULL PLAYER TABLE, GROUPED BY TEAM, WITH
#    CONDITIONAL FORMATTING (highest/lowest per column highlighted)
# ============================================================

def add_appendix_table_slide(prs, section_title, phase_name, team_name, team_color,
                              player_df, metrics):
    """
    Full roster table (every player) for the metrics in `metrics`
    (list of (label, key, direction) tuples). The best value in each
    column is highlighted green, the worst is highlighted red -
    "conditional formatting" applied at build time since PowerPoint
    tables don't support live conditional formatting rules.
    """
    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_LIGHT)

    _add_textbox(slide, Inches(0.5), Inches(0.28), Inches(9.0), Inches(0.3),
                 section_title.upper(), size=11, bold=True, color=MUTED_TEXT, letter_spaced=True)
    _add_textbox(slide, Inches(0.5), Inches(0.55), Inches(9.0), Inches(0.55),
                 f"{team_name} \u2014 {phase_name}", size=22, bold=True, color=DARK_TEXT, font_name=HEADER_FONT)

    if player_df is None or player_df.empty or "Player" not in player_df.columns:
        _add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.5),
                     "No player data available", size=13, color=MUTED_TEXT)
        _add_footer(slide, prs, f"Appendix \u2014 {phase_name}", team_name.upper())
        return slide

    available_metrics = [(label, key, direction) for label, key, direction in metrics
                          if key in player_df.columns]

    df = player_df[["Player"] + [key for _, key, _ in available_metrics]].fillna(0).copy()
    df = df.sort_values("Player").reset_index(drop=True)

    n_rows = len(df) + 1
    n_cols = len(available_metrics) + 1

    table_top = Inches(1.3)
    table_left = Inches(0.5)
    table_width = Inches(SLIDE_WIDTH_IN - 1.0)
    table_height = Inches(min(5.6, 0.32 * n_rows + 0.3))

    gshape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
    table = gshape.table

    table.columns[0].width = Emu(int(table_width * 0.28))
    remaining_width = int(table_width * 0.72 / max(len(available_metrics), 1))
    for c in range(1, n_cols):
        table.columns[c].width = Emu(remaining_width)

    table.cell(0, 0).text = "Player"
    for c, (label, _, _) in enumerate(available_metrics, start=1):
        table.cell(0, c).text = label

    for r, row in enumerate(df.itertuples(index=False), start=1):
        table.cell(r, 0).text = str(row.Player)
        for c, (_, key, _) in enumerate(available_metrics, start=1):
            val = getattr(row, key)
            table.cell(r, c).text = str(int(val)) if float(val).is_integer() else f"{val:.2f}"
            table.cell(r, c).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    _style_table(table, team_color, _readable_text_color(team_color))

    # Conditional formatting: highlight best/worst value per metric column
    for c, (_, key, direction) in enumerate(available_metrics, start=1):
        col_values = df[key]
        if col_values.empty:
            continue
        max_val, min_val = col_values.max(), col_values.min()
        if max_val == min_val:
            continue
        best_val = max_val if direction != "lower" else min_val
        worst_val = min_val if direction != "lower" else max_val
        for r, val in enumerate(col_values, start=1):
            cell = table.cell(r, c)
            if val == best_val:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HIGHLIGHT_BEST
            elif val == worst_val:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HIGHLIGHT_WORST

    _add_footer(slide, prs, f"Appendix \u2014 {phase_name}", team_name.upper())
    return slide


# ============================================================
# 6. AVERAGE POSITION / SHOT CHART IMAGE SLIDE
# ============================================================

def add_image_slide(prs, title, image_path, team_name, team_color, caption=None):
    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_LIGHT)
    _add_textbox(slide, Inches(0.5), Inches(0.35), Inches(SLIDE_WIDTH_IN - 1.0), Inches(0.6),
                 title, size=26, bold=True, color=DARK_TEXT, font_name=HEADER_FONT)

    max_w = Inches(SLIDE_WIDTH_IN - 1.6)
    max_h = Inches(5.1)
    top = Inches(1.2)

    with Image.open(image_path) as im:
        img_w, img_h = im.size
    aspect = img_w / img_h

    if (max_w / max_h) > aspect:
        height = max_h
        width = Emu(int(height * aspect))
    else:
        width = max_w
        height = Emu(int(width / aspect))

    left = Emu(int((prs.slide_width - width) / 2))
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

    if caption:
        _add_textbox(slide, Inches(0.5), Inches(6.5), Inches(SLIDE_WIDTH_IN - 1.0), Inches(0.4),
                     caption, size=12, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

    _add_footer(slide, prs, title, team_name.upper())
    return slide


# ============================================================
# 7. SUPPLEMENTARY KPI TABLE
# ============================================================

def add_supplementary_table_slide(prs, comparison_df, team_name, opponent_name, team_color, opponent_color,
                                   metrics=None, title="Additional KPIs"):
    metrics = metrics or sp.SUPPLEMENTARY_METRICS
    labels = [label for label, _, _ in metrics]
    rows = comparison_df[comparison_df["Metric"].isin(labels)]
    rows = rows.set_index("Metric").reindex(labels).reset_index()

    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_LIGHT)
    _add_textbox(slide, Inches(0.5), Inches(0.35), Inches(SLIDE_WIDTH_IN - 1.0), Inches(0.6),
                 title, size=26, bold=True, color=DARK_TEXT, font_name=HEADER_FONT)

    n_rows = len(rows) + 1
    table_top = Inches(1.4)
    table_left = Inches(2.7)
    table_width = Inches(SLIDE_WIDTH_IN - 5.4)
    table_height = Inches(min(4.8, 0.6 * n_rows))

    gshape = slide.shapes.add_table(n_rows, 3, table_left, table_top, table_width, table_height)
    table = gshape.table
    table.columns[0].width = Emu(int(table_width * 0.42))
    table.columns[1].width = Emu(int(table_width * 0.29))
    table.columns[2].width = Emu(int(table_width * 0.29))

    table.cell(0, 0).text = "KPI"
    table.cell(0, 1).text = team_name
    table.cell(0, 2).text = opponent_name

    for i, record in enumerate(rows.to_dict(orient="records"), start=1):
        table.cell(i, 0).text = str(record["Metric"])
        val_team = record[team_name]
        val_opp = record[opponent_name]
        table.cell(i, 1).text = "-" if pd.isna(val_team) else str(val_team)
        table.cell(i, 2).text = "-" if pd.isna(val_opp) else str(val_opp)
        for c in range(3):
            table.cell(i, c).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    _style_table(table, team_color, _readable_text_color(team_color))
    _add_footer(slide, prs, "Supplementary Statistics", f"{team_name.upper()} vs {opponent_name.upper()}")
    return slide


# ============================================================
# 8. FULL-TIME SUMMARY SLIDE (direction-aware scoring)
# ============================================================

def add_fulltime_slide(prs, comparison_df, team_name, opponent_name, team_color, opponent_color,
                        metrics=None):
    metrics = metrics or sp.FLAT_MAIN_METRICS
    rows = comparison_df[comparison_df["Metric"].isin([label for label, _, _ in metrics])]
    direction_by_label = {label: direction for label, _, direction in metrics}

    team_wins = 0
    opponent_wins = 0
    ties = 0

    for record in rows.to_dict(orient="records"):
        team_val = record[team_name]
        opponent_val = record[opponent_name]
        direction = direction_by_label.get(record["Metric"], "higher")

        if pd.isna(team_val) or pd.isna(opponent_val) or team_val == opponent_val:
            ties += 1
        elif (direction == "lower" and team_val < opponent_val) or \
             (direction != "lower" and team_val > opponent_val):
            team_wins += 1
        else:
            opponent_wins += 1

    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_DARK)
    _add_decorative_arc(slide, prs)

    _add_textbox(slide, Inches(0.8), Inches(0.6), Inches(6.0), Inches(0.4),
                 "F U L L - T I M E", size=14, color=GREY_LABEL)

    if team_wins > opponent_wins:
        headline = f"{team_name} leads the match statistics"
    elif opponent_wins > team_wins:
        headline = f"{opponent_name} leads the match statistics"
    else:
        headline = "Match statistics are evenly split"

    _add_textbox(slide, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.7),
                 headline, size=30, bold=True, color=CREAM, font_name=HEADER_FONT)

    opp_luminance = 0.299 * opponent_color[0] + 0.587 * opponent_color[1] + 0.114 * opponent_color[2]
    opp_display_color = CREAM if opp_luminance < 60 else opponent_color

    stat_specs = [
        (team_wins, f"categories won\nby {team_name}", team_color),
        (opponent_wins, f"categories won\nby {opponent_name}", opp_display_color),
        (ties, "categories tied", GREY_LABEL),
    ]

    box_w = Inches(3.4)
    gap = Inches(0.6)
    total_w = box_w * 3 + gap * 2
    start_x = Emu(int((prs.slide_width - total_w) / 2))
    top = Inches(2.9)

    for i, (value, caption, color) in enumerate(stat_specs):
        left = Emu(int(start_x + i * (box_w + gap)))
        _add_textbox(slide, left, top, box_w, Inches(1.3),
                     str(value), size=64, bold=True, color=color, align=PP_ALIGN.CENTER,
                     font_name=HEADER_FONT)
        cap_box = slide.shapes.add_textbox(left, top + Inches(1.35), box_w, Inches(0.7))
        tf = cap_box.text_frame
        tf.word_wrap = True
        for j, line in enumerate(caption.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            run.font.color.rgb = GREY_LABEL
            run.font.name = BODY_FONT

    _add_footer(slide, prs, "Match Recap Complete", f"{team_name.upper()} vs {opponent_name.upper()}", on_dark=True)
    return slide


# ============================================================
# MAIN ENTRY POINT
# ============================================================

def generate_report(
    comparison_csv,
    team_name,
    team_color,
    opponent_name,
    opponent_color,
    output_path,
    team_score=None,
    opponent_score=None,
    date_label=None,
    kicker="STATSBOMB MATCH RECAP",
    subtitle="STATISTICAL MATCH ANALYSIS  \u2022  PLAYER & TEAM PERFORMANCE BREAKDOWN",
    team_player_stats_csv=None,
    opponent_player_stats_csv=None,
    team_avg_position_image=None,
    opponent_avg_position_image=None,
    team_shot_chart_image=None,
    opponent_shot_chart_image=None,
    logo_image=None,
    include_supplementary_table=True,
    include_goalkeepers=True,
):
    team_rgb = _hex_to_rgb(team_color)
    opponent_rgb = _hex_to_rgb(opponent_color)

    comparison_df = pd.read_csv(comparison_csv)
    for col in (team_name, opponent_name):
        comparison_df[col] = pd.to_numeric(comparison_df[col], errors="coerce")

    team_players_df = pd.read_csv(team_player_stats_csv) if team_player_stats_csv is not None else None
    opponent_players_df = pd.read_csv(opponent_player_stats_csv) if opponent_player_stats_csv is not None else None

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    # 1. Title slide
    add_title_slide(prs, team_name, team_rgb, opponent_name, opponent_rgb,
                     kicker=kicker, subtitle=subtitle, date_label=date_label,
                     logo_image=logo_image)

    # Which phase each pair of team visuals logically belongs to
    avg_position_phase = "Build-Up & Possession"
    shot_chart_phase = "Final Third / Red Zone"

    # 2. Game-phase sections: divider, comparison, per-metric slides,
    #    appendix, and any visuals that belong to this phase
    phase_count = len(sp.GAME_PHASES)

    for phase_index, phase in enumerate(sp.GAME_PHASES, start=1):
        phase_name = phase["name"]
        phase_metrics = phase["metrics"]

        add_phase_title_slide(
            prs, phase_index, phase_count, phase_name, phase_metrics,
            team_name, opponent_name, team_rgb, opponent_rgb,
            logo_image=logo_image,
        )

        add_phase_comparison_slide(
            prs, phase_name, comparison_df, phase_metrics,
            team_name, opponent_name, team_rgb, opponent_rgb,
        )

        player_metrics = [m for m in phase_metrics if m[1] != "ppda"]
        stat_total = len(player_metrics)
        phase_rows = comparison_df.set_index("Metric").reindex([label for label, _, _ in player_metrics])

        for i, (label, player_col, direction) in enumerate(player_metrics, start=1):
            team_value = phase_rows.loc[label, team_name] if label in phase_rows.index else None
            opponent_value = phase_rows.loc[label, opponent_name] if label in phase_rows.index else None
            add_metric_slide(
                prs, phase_name, label, player_col, team_value, opponent_value,
                team_players_df, opponent_players_df,
                team_name, opponent_name, team_rgb, opponent_rgb,
                direction=direction, include_goalkeepers=include_goalkeepers,
                stat_index=i, stat_total=stat_total,
            )

        # Appendix for this phase - full roster, conditional formatting
        if player_metrics:
            add_appendix_table_slide(prs, "Appendix", phase_name, team_name, team_rgb,
                                      team_players_df, player_metrics)
            add_appendix_table_slide(prs, "Appendix", phase_name, opponent_name, opponent_rgb,
                                      opponent_players_df, player_metrics)

        # Team visuals that logically belong to this phase, grouped by team
        if phase_name == avg_position_phase:
            if team_avg_position_image is not None:
                add_image_slide(prs, f"{team_name} \u2014 Average Position", team_avg_position_image, team_name, team_rgb)
            if opponent_avg_position_image is not None:
                add_image_slide(prs, f"{opponent_name} \u2014 Average Position", opponent_avg_position_image, opponent_name, opponent_rgb)

        if phase_name == shot_chart_phase:
            if team_shot_chart_image is not None:
                add_image_slide(prs, f"{team_name} \u2014 Shot Chart", team_shot_chart_image, team_name, team_rgb)
            if opponent_shot_chart_image is not None:
                add_image_slide(prs, f"{opponent_name} \u2014 Shot Chart", opponent_shot_chart_image, opponent_name, opponent_rgb)

    # 3. Supplementary KPI table
    if include_supplementary_table:
        add_supplementary_table_slide(prs, comparison_df, team_name, opponent_name, team_rgb, opponent_rgb)

    # 4. Full-time summary slide
    add_fulltime_slide(prs, comparison_df, team_name, opponent_name, team_rgb, opponent_rgb)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return str(output_path)
